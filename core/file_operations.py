import os
import re
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation


#text cleaning//

def clean_text(text):
    """Basic cleaning to normalize extracted text."""
    if not text:
        return ""

    text = re.sub(r"\s+", " ", text)
    return text.strip()



#pdf recusive scan//

def list_pdf_files_recursively(folder_path):
    """
    Recursively scans a folder and returns all PDF file paths.
    """
    pdf_files = []

    for root, dirs, files in os.walk(folder_path):
        for file in files:
            if file.lower().endswith(".pdf"):
                pdf_files.append(os.path.join(root, file))

    return pdf_files



#pdf text extraction//

def extract_text_from_pdf(pdf_path):
    """Extract text from a PDF file."""
    try:
        reader = PdfReader(pdf_path)
        text = ""

        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"

        return clean_text(text)

    except Exception as e:
        print(f"Error extracting PDF ({pdf_path}): {e}")
        return ""



#docx text extraction//

def extract_text_from_docx(docx_path):
    """Extract text from a DOCX file."""
    try:
        doc = Document(docx_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return clean_text(text)

    except Exception as e:
        print(f"Error extracting DOCX ({docx_path}): {e}")
        return ""



#pptx text extraction//

def extract_text_from_pptx(pptx_path):
    """Extract text from PPTX slides."""
    try:
        prs = Presentation(pptx_path)
        text = ""

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

        return clean_text(text)

    except Exception as e:
        print(f"Error extracting PPTX ({pptx_path}): {e}")
        return ""


# Extract all pdfs in folder//

def extract_pdfs_into_list(folder_path):
    """
    Returns list of dicts: [{path: "", text: ""}, ...]
    Recursively scans all subfolders.
    """
    pdf_files = list_pdf_files_recursively(folder_path)
    output = []

    for pdf in pdf_files:
        text = extract_text_from_pdf(pdf)
        output.append({"path": pdf, "text": text})

    return output
